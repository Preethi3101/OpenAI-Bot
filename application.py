import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.question_answering import load_qa_chain
from langchain.chat_models import ChatOpenAI
from langchain.prompts import PromptTemplate
import os
from pptx import Presentation
from docx import Document
import zipfile
from io import BytesIO
from azure.cognitiveservices.vision.computervision import ComputerVisionClient
from azure.cognitiveservices.vision.computervision.models import OperationStatusCodes
from msrest.authentication import CognitiveServicesCredentials

# Load environment variables
load_dotenv()
cv_key = os.getenv("AZURE_VISION_KEY")
cv_endpoint = os.getenv("AZURE_VISION_ENDPOINT")

# Initialize Azure Computer Vision client
cv_client = ComputerVisionClient(cv_endpoint, CognitiveServicesCredentials(cv_key))

def get_image_text(image):
    read_response = cv_client.read_in_stream(image, raw=True)
    read_operation_location = read_response.headers["Operation-Location"]
    operation_id = read_operation_location.split("/")[-1]

    while True:
        read_result = cv_client.get_read_result(operation_id)
        if read_result.status not in [OperationStatusCodes.not_started, OperationStatusCodes.running]:
            break
    text = ""
    if read_result.status == OperationStatusCodes.succeeded:
        for page in read_result.analyze_result.read_results:
            for line in page.lines:
                text += line.text + "\n"
    return text

def get_pdf_text(pdf):
    text = ""
    pdf_reader = PdfReader(pdf)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_ppt_text(ppt):
    text = ""
    presentation = Presentation(ppt)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def get_docx_text(docx):
    text = ""
    document = Document(docx)
    for para in document.paragraphs:
        text += para.text
    return text

def get_text_from_file(file):
    filename = file.name
    if filename.endswith('.pdf'):
        return get_pdf_text(file)
    elif filename.endswith('.pptx'):
        return get_ppt_text(file)
    elif filename.endswith('.docx'):
        return get_docx_text(file)
    elif filename.lower().endswith(('.png', '.jpg', '.jpeg')):
        return get_image_text(file)
    else:
        return ""

def get_text_from_zip(zip_file):
    text = ""
    with zipfile.ZipFile(BytesIO(zip_file.read()), "r") as z:
        for file_name in z.namelist():
            if file_name.lower().endswith(('.pdf', '.pptx', '.docx', '.png', '.jpg', '.jpeg')):
                with z.open(file_name) as inner_file:
                    text += get_text_from_file(inner_file)
    return text

def get_text_chunks(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = splitter.split_text(text)
    return chunks

def get_vectorstore(chunks, embeddings):
    vectorstore = FAISS.from_texts(texts=chunks, embedding=embeddings)
    vectorstore.save_local("faiss_index")

def clear_chat_history():
    st.session_state.messages = [{"role": "assistant", "content": "Upload some documents and ask me a question."}]

def get_conversation_chain():
    prompt_template = """
    You are a professional in understanding content from PDFs and answering any questions related to it. You are also a helpful assistant who knows how to greet people. Understand the content provided and answer the questions appropriately. Ensure that the sentences you provide are grammatically correct. Answer the question as detailed as possible from the provided context for all questions except greetings; make sure to provide all the details. If the answer is not in
    the provided context, just say, "The answer is not available in the context." Don't provide the wrong answer. If it is greetings, please answer appropriately.\n\n
    Context:\n {context}?\n
    Question: \n{question}\n
    Answer:
    """
    model = ChatOpenAI()
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(llm=model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question, embeddings):
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversation_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response

def main():
    st.set_page_config(page_title="OpenAI PDF/PPT Chatbot", page_icon="book")

    st.title("Chat with Documents")
    st.write("Welcome to the chat!")

    with st.sidebar:
        st.title("Menu:")
        zip_file = st.file_uploader("Upload your ZIP File", type="zip")
        if st.button("Submit & Process"):
            if zip_file is not None:
                text = get_text_from_zip(zip_file)
                if text:
                    chunks = get_text_chunks(text)
                    embeddings = OpenAIEmbeddings()
                    get_vectorstore(chunks, embeddings)
                    st.success("Text processed and vectorstore created successfully.")
                else:
                    st.warning("No valid documents found in the ZIP file.")
            else:
                st.warning("Please upload a ZIP file.")

    st.sidebar.button('Clear Chat History', on_click=clear_chat_history)

    if "messages" not in st.session_state.keys():
        st.session_state.messages = [{"role": "assistant", "content": "Upload some documents and ask me a question."}]

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.write(message["content"])

    if prompt := st.chat_input():
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.write(prompt)

        if st.session_state.messages[-1]["role"] != "assistant":
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    embeddings = OpenAIEmbeddings()
                    response = user_input(prompt, embeddings)
                    placeholder = st.empty()
                    full_response = ''
                    for item in response['output_text']:
                        full_response += item
                        placeholder.markdown(full_response)
                    placeholder.markdown(full_response)

            if response is not None:
                message = {"role": "assistant", "content": full_response}
                st.session_state.messages.append(message)

if __name__ == "__main__":
    main()
